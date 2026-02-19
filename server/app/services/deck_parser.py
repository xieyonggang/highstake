import io
import json
import logging
import tempfile
import uuid
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


class DeckParserService:
    async def parse_and_store(
        self,
        file_bytes: bytes,
        filename: str,
        session_id: str | None = None,
    ) -> dict:
        """Full pipeline: parse file, generate thumbnails, write manifest.json.

        If session_id is provided, files are stored directly in the session
        folder at sessions/{session_id}/decks/{deck_id}/ instead of decks/{deck_id}/.
        """
        filename_lower = filename.lower()
        is_pptx = filename_lower.endswith(".pptx")
        is_pdf = filename_lower.endswith(".pdf")

        if not (is_pptx or is_pdf):
            raise ValueError("Unsupported file format. Only PPTX and PDF are supported.")

        with tempfile.TemporaryDirectory() as tmpdir:
            tmpdir_path = Path(tmpdir)
            file_path = tmpdir_path / filename
            file_path.write_bytes(file_bytes)

            if is_pptx:
                slides_data = self._parse_pptx(file_path)
            else:
                slides_data = self._parse_pdf(file_path)

            # Generate thumbnails
            thumbnails = self._generate_thumbnails(file_path, tmpdir_path, is_pptx)

            # Storage prefix: session folder or standalone decks folder
            from app.services.storage_service import StorageService

            storage = StorageService()
            deck_id = str(uuid.uuid4())
            if session_id:
                prefix = f"sessions/{session_id}/decks/{deck_id}"
            else:
                prefix = f"decks/{deck_id}"

            file_key = f"{prefix}/{filename}"
            await storage.upload(file_key, file_bytes, self._content_type(filename))

            # Upload thumbnails
            for i, thumb_bytes in enumerate(thumbnails):
                if thumb_bytes:
                    thumb_key = f"{prefix}/thumbnails/{i}.png"
                    await storage.upload(thumb_key, thumb_bytes, "image/png")
                    slides_data[i]["thumbnail_key"] = thumb_key

            # Save parsed slide content as readable markdown
            slides_md = self._build_slides_markdown(filename, slides_data)
            slides_md_key = f"{prefix}/slides.md"
            await storage.upload(
                slides_md_key,
                slides_md.encode("utf-8"),
                "text/markdown",
            )

            # Build manifest with direct file URLs for thumbnails
            slides_out = []
            for s in slides_data:
                slides_out.append({
                    "index": s["index"],
                    "title": s.get("title"),
                    "subtitle": s.get("subtitle"),
                    "body_text": s.get("body_text"),
                    "notes": s.get("notes"),
                    "has_chart": s.get("has_chart", False),
                    "has_table": s.get("has_table", False),
                    "thumbnail_url": (
                        f"/api/files/{prefix}/thumbnails/{s['index']}.png"
                        if s.get("thumbnail_key")
                        else None
                    ),
                })

            manifest = {
                "id": deck_id,
                "filename": filename,
                "totalSlides": len(slides_data),
                "slides": slides_out,
            }

            # Write manifest.json to storage
            manifest_key = f"{prefix}/manifest.json"
            await storage.upload(
                manifest_key,
                json.dumps(manifest, indent=2).encode(),
                "application/json",
            )

            return {
                "id": deck_id,
                "filename": filename,
                "total_slides": len(slides_data),
                "slides": slides_out,
            }

    def _parse_pptx(self, file_path: Path) -> list[dict]:
        """Extract text, titles, notes from PPTX."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(file_path))
        slides_data = []

        for i, slide in enumerate(prs.slides):
            title = ""
            subtitle = ""
            body_parts = []
            has_chart = False
            has_table = False
            notes = ""

            # Extract title
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()

            for shape in slide.shapes:
                # Check shape types
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    has_chart = True
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                    # Extract table text
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                body_parts.append(row_text)

                # Extract text from text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            if shape == slide.shapes.title:
                                continue  # Already captured
                            elif not subtitle and shape.shape_id != slide.shapes.title.shape_id if slide.shapes.title else True:
                                # First non-title text block could be subtitle
                                if not subtitle and len(body_parts) == 0:
                                    subtitle = text
                                else:
                                    body_parts.append(text)

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes = notes_frame.text.strip()

            slides_data.append({
                "index": i,
                "title": title,
                "subtitle": subtitle,
                "body_text": "\n".join(body_parts),
                "notes": notes,
                "has_chart": has_chart,
                "has_table": has_table,
            })

        return slides_data

    def _parse_pdf(self, file_path: Path) -> list[dict]:
        """Extract text from PDF pages."""
        import fitz  # PyMuPDF

        doc = fitz.open(str(file_path))
        slides_data = []

        for i, page in enumerate(doc):
            text = page.get_text("text").strip()
            lines = [line.strip() for line in text.split("\n") if line.strip()]

            title = lines[0] if lines else f"Page {i + 1}"
            subtitle = lines[1] if len(lines) > 1 else ""
            body_text = "\n".join(lines[2:]) if len(lines) > 2 else ""

            slides_data.append({
                "index": i,
                "title": title,
                "subtitle": subtitle,
                "body_text": body_text,
                "notes": "",
                "has_chart": False,
                "has_table": False,
            })

        doc.close()
        return slides_data

    def _generate_thumbnails(
        self,
        file_path: Path,
        output_dir: Path,
        is_pptx: bool,
    ) -> list[Optional[bytes]]:
        """Generate thumbnail images for each slide/page."""
        thumbnails = []

        if is_pptx:
            # For PPTX, we generate text-based placeholder thumbnails
            # Full PPTX rendering requires LibreOffice or similar
            try:
                from pptx import Presentation
                from PIL import Image, ImageDraw, ImageFont

                prs = Presentation(str(file_path))
                for i, slide in enumerate(prs.slides):
                    thumb = self._create_text_thumbnail(
                        title=slide.shapes.title.text.strip() if slide.shapes.title else f"Slide {i + 1}",
                        index=i,
                    )
                    thumbnails.append(thumb)
            except Exception as e:
                logger.warning(f"Failed to generate PPTX thumbnails: {e}")

        else:
            # For PDF, use PyMuPDF to render pages
            try:
                import fitz

                doc = fitz.open(str(file_path))
                for page in doc:
                    mat = fitz.Matrix(2, 2)  # 2x zoom
                    pix = page.get_pixmap(matrix=mat)
                    thumbnails.append(pix.tobytes("png"))
                doc.close()
            except Exception as e:
                logger.warning(f"Failed to generate PDF thumbnails: {e}")

        return thumbnails

    def _create_text_thumbnail(self, title: str, index: int) -> bytes:
        """Create a simple text-based thumbnail for a PPTX slide."""
        from PIL import Image, ImageDraw, ImageFont

        width, height = 1280, 720
        img = Image.new("RGB", (width, height), color=(30, 30, 40))
        draw = ImageDraw.Draw(img)

        # Draw slide number
        draw.text((40, 30), f"Slide {index + 1}", fill=(100, 100, 140))

        # Draw title
        # Use default font (PIL built-in)
        title_wrapped = title[:60]
        draw.text((40, 80), title_wrapped, fill=(220, 220, 240))

        # Draw border
        draw.rectangle([0, 0, width - 1, height - 1], outline=(60, 60, 80), width=2)

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    def _build_slides_markdown(self, filename: str, slides_data: list[dict]) -> str:
        """Build a human-readable markdown summary of all parsed slides."""
        lines = [
            f"# {filename}",
            f"",
            f"**Total slides:** {len(slides_data)}",
            "",
        ]
        for s in slides_data:
            idx = s.get("index", 0)
            title = s.get("title", f"Slide {idx + 1}")
            lines.append(f"---\n## Slide {idx + 1}: {title}\n")
            if s.get("subtitle"):
                lines.append(f"**Subtitle:** {s['subtitle']}\n")
            if s.get("body_text"):
                lines.append(f"{s['body_text']}\n")
            if s.get("notes"):
                lines.append(f"**Speaker notes:** {s['notes']}\n")
            flags = []
            if s.get("has_chart"):
                flags.append("chart")
            if s.get("has_table"):
                flags.append("table")
            if flags:
                lines.append(f"_Contains: {', '.join(flags)}_\n")
        return "\n".join(lines)

    def _content_type(self, filename: str) -> str:
        if filename.lower().endswith(".pptx"):
            return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif filename.lower().endswith(".pdf"):
            return "application/pdf"
        return "application/octet-stream"
